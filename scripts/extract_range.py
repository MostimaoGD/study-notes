#!/usr/bin/env python
"""Extract selected learning-document ranges with source anchors."""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Iterable


RANGE_RE = re.compile(r"^\s*(\d+)(?:\s*-\s*(\d+))?\s*$")


def parse_range_spec(spec: str | None, max_value: int | None = None) -> list[int] | None:
    if not spec:
        return None
    values: set[int] = set()
    for part in spec.split(","):
        match = RANGE_RE.match(part)
        if not match:
            raise ValueError(f"Invalid range segment: {part!r}")
        start = int(match.group(1))
        end = int(match.group(2) or start)
        if end < start:
            start, end = end, start
        for value in range(start, end + 1):
            if value >= 1 and (max_value is None or value <= max_value):
                values.add(value)
    return sorted(values)


def read_text_file(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "utf-16"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def split_plain_blocks(text: str, max_block_chars: int = 6000) -> list[str]:
    blocks: list[str] = []
    for part in re.split(r"\n\s*\n", text):
        part = part.strip()
        if not part:
            continue
        if len(part) <= max_block_chars:
            blocks.append(part)
            continue
        lines = [line.strip() for line in part.splitlines() if line.strip()]
        if len(lines) > 1 and max(len(line) for line in lines) <= max_block_chars:
            blocks.extend(lines)
            continue
        for start in range(0, len(part), max_block_chars):
            blocks.append(part[start : start + max_block_chars].strip())
    return [block for block in blocks if block]


def clip(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + "\n\n<!-- extraction truncated by --max-chars -->\n"


def extract_pdf(path: Path, pages_spec: str | None, max_chars: int) -> str:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        pages = parse_range_spec(pages_spec, len(pdf.pages)) or list(range(1, len(pdf.pages) + 1))
        for page_number in pages:
            page = pdf.pages[page_number - 1]
            text = page.extract_text() or ""
            parts.append(f"<!-- source: {path.name} p.{page_number} -->\n\n{text.strip()}")
    return clip("\n\n".join(parts).strip() + "\n", max_chars)


def iter_shape_text(shape) -> list[str]:
    texts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip()
        if text:
            texts.append(text)
    if getattr(shape, "has_table", False):
        try:
            rows = []
            for row in shape.table.rows:
                rows.append(" | ".join(cell.text.strip() for cell in row.cells))
            if rows:
                texts.append("\n".join(rows))
        except Exception:
            pass
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            texts.extend(iter_shape_text(child))
    return texts


def extract_pptx(path: Path, slides_spec: str | None, max_chars: int) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    slides = parse_range_spec(slides_spec, len(prs.slides)) or list(range(1, len(prs.slides) + 1))
    parts: list[str] = []
    for slide_number in slides:
        slide = prs.slides[slide_number - 1]
        texts: list[str] = []
        for shape in slide.shapes:
            texts.extend(iter_shape_text(shape))
        parts.append(f"<!-- source: {path.name} slide {slide_number} -->\n\n" + "\n\n".join(texts).strip())
    return clip("\n\n".join(parts).strip() + "\n", max_chars)


def extract_docx(path: Path, paragraphs_spec: str | None, heading: str | None, max_chars: int) -> str:
    from docx import Document

    doc = Document(str(path))
    paragraphs = list(doc.paragraphs)
    selected_indexes: Iterable[int]
    if heading:
        start = None
        start_level = None
        heading_lower = heading.lower()
        for index, paragraph in enumerate(paragraphs):
            style_name = (paragraph.style.name if paragraph.style else "").lower()
            if heading_lower in paragraph.text.lower():
                start = index
                match = re.search(r"heading\s+(\d+)", style_name)
                start_level = int(match.group(1)) if match else 9
                break
        if start is None:
            return f"<!-- source: {path.name} heading not found: {heading} -->\n"
        end = len(paragraphs)
        for index in range(start + 1, len(paragraphs)):
            style_name = (paragraphs[index].style.name if paragraphs[index].style else "").lower()
            match = re.search(r"heading\s+(\d+)", style_name)
            if match and int(match.group(1)) <= int(start_level or 9):
                end = index
                break
        selected_indexes = range(start, end)
    else:
        selected = parse_range_spec(paragraphs_spec, len(paragraphs))
        selected_indexes = [index - 1 for index in selected] if selected else range(len(paragraphs))

    parts = []
    for index in selected_indexes:
        text = paragraphs[index].text.strip()
        if text:
            parts.append(f"<!-- source: {path.name} paragraph {index + 1} -->\n\n{text}")
    return clip("\n\n".join(parts).strip() + "\n", max_chars)


def markdown_heading_level(line: str) -> int | None:
    match = re.match(r"^(#{1,6})\s+\S", line)
    return len(match.group(1)) if match else None


def extract_markdown_or_txt(path: Path, paragraphs_spec: str | None, heading: str | None, max_chars: int) -> str:
    text = read_text_file(path)
    if heading:
        lines = text.splitlines()
        heading_lower = heading.lower()
        start = None
        start_level = None
        for index, line in enumerate(lines):
            if heading_lower in line.lower():
                start = index
                start_level = markdown_heading_level(line) or 6
                break
        if start is None:
            return f"<!-- source: {path.name} heading not found: {heading} -->\n"
        end = len(lines)
        for index in range(start + 1, len(lines)):
            level = markdown_heading_level(lines[index])
            if level is not None and level <= int(start_level or 6):
                end = index
                break
        return clip(f"<!-- source: {path.name} heading {heading} -->\n\n" + "\n".join(lines[start:end]).strip() + "\n", max_chars)

    paragraphs = split_plain_blocks(text)
    selected = parse_range_spec(paragraphs_spec, len(paragraphs))
    if selected:
        parts = [
            f"<!-- source: {path.name} paragraph {number} -->\n\n{paragraphs[number - 1].strip()}"
            for number in selected
            if paragraphs[number - 1].strip()
        ]
    else:
        parts = [f"<!-- source: {path.name} -->\n\n{text.strip()}"]
    return clip("\n\n".join(parts).strip() + "\n", max_chars)


def extract_range(path: Path, args: argparse.Namespace) -> str:
    if args.plan_json:
        if getattr(args, "all_items", False):
            return extract_plan_items(path, args.plan_json, args.max_chars)
        return extract_plan_batch(path, args.plan_json, args.batch_id or "first", args.max_chars)
    ext = path.suffix.lower()
    if ext == ".pdf":
        return extract_pdf(path, args.pages, args.max_chars)
    if ext == ".pptx":
        return extract_pptx(path, args.slides, args.max_chars)
    if ext == ".docx":
        return extract_docx(path, args.paragraphs, args.heading, args.max_chars)
    if ext in {".md", ".markdown", ".txt"}:
        return extract_markdown_or_txt(path, args.paragraphs, args.heading, args.max_chars)
    raise ValueError(f"Unsupported file extension: {ext}")


def select_batch(plan: dict, batch_id: str) -> dict:
    selected_id = plan.get("first_batch") if batch_id == "first" else batch_id
    for batch in plan.get("batches", []):
        if batch.get("id") == selected_id:
            return batch
    raise ValueError(f"Batch not found: {batch_id}")


def load_plan(plan_json: str) -> dict:
    return json.loads(Path(plan_json).read_text(encoding="utf-8"))


def warn_if_plan_file_differs(path: Path, plan: dict) -> None:
    planned_file = plan.get("file")
    if planned_file and Path(planned_file).name != path.name:
        sys.stderr.write(f"Warning: plan file name differs from input file: {planned_file}\n")


def item_scope(item: dict) -> str:
    extraction = item.get("extraction") or {}
    if "pages" in extraction:
        return f"p.{extraction['pages']}"
    if "slides" in extraction:
        return f"slide {extraction['slides']}"
    if "paragraphs" in extraction:
        return f"paragraph {extraction['paragraphs']}"
    if "heading" in extraction:
        return f"heading {extraction['heading']}"
    return "unknown"


def markdown_cell(value: object) -> str:
    return str(value or "").replace("|", "\\|").replace("\n", " ").strip()


def slugify_filename(text: object, fallback: str, limit: int = 56) -> str:
    cleaned = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', "-", str(text or ""))
    cleaned = re.sub(r"\s+", " ", cleaned).strip(" .-_")
    if not cleaned:
        cleaned = fallback
    return (cleaned[:limit].strip(" .-_") or fallback)


def extract_item(path: Path, item: dict, max_chars: int) -> str:
    extraction = item.get("extraction") or {}
    title = item.get("title") or item.get("id") or "planned item"
    header = f"\n\n<!-- planned-item: {item.get('id', '')} | {title} -->\n\n"
    if "pages" in extraction:
        return header + extract_pdf(path, str(extraction["pages"]), max_chars)
    if "slides" in extraction:
        return header + extract_pptx(path, str(extraction["slides"]), max_chars)
    if "paragraphs" in extraction:
        ext = path.suffix.lower()
        if ext == ".docx":
            return header + extract_docx(path, str(extraction["paragraphs"]), None, max_chars)
        if ext in {".md", ".markdown", ".txt"}:
            return header + extract_markdown_or_txt(path, str(extraction["paragraphs"]), None, max_chars)
    if "heading" in extraction:
        ext = path.suffix.lower()
        if ext == ".docx":
            return header + extract_docx(path, None, str(extraction["heading"]), max_chars)
        if ext in {".md", ".markdown", ".txt"}:
            return header + extract_markdown_or_txt(path, None, str(extraction["heading"]), max_chars)
    raise ValueError(f"Unsupported extraction item: {item}")


def extract_plan_batch(path: Path, plan_json: str, batch_id: str, max_chars: int) -> str:
    plan = load_plan(plan_json)
    batch = select_batch(plan, batch_id)
    warn_if_plan_file_differs(path, plan)
    parts = [
        f"<!-- planned-batch: {batch.get('id')} -->",
        f"<!-- structure-source: {plan.get('structure_source')} | confidence: {plan.get('confidence')} -->",
        "",
    ]
    for item in batch.get("items", []):
        parts.append(extract_item(path, item, max_chars))
    return clip("\n\n".join(parts).strip() + "\n", max_chars)


def extract_plan_items(path: Path, plan_json: str, max_chars: int) -> str:
    plan = load_plan(plan_json)
    warn_if_plan_file_differs(path, plan)
    parts = [
        "<!-- planned-extraction: all-items -->",
        f"<!-- structure-source: {plan.get('structure_source')} | confidence: {plan.get('confidence')} -->",
        "",
    ]
    for item in plan.get("chapters", []):
        parts.append(extract_item(path, item, max_chars))
    return clip("\n\n".join(parts).strip() + "\n", max_chars)


def extract_plan_items_to_dir(path: Path, plan_json: str, output_dir: str, max_chars: int) -> str:
    plan = load_plan(plan_json)
    warn_if_plan_file_differs(path, plan)
    out_dir = Path(output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    items = list(plan.get("chapters", []))
    index_lines = [
        "# 全部章节/分块抽取索引",
        "",
        f"- 来源文件：`{plan.get('file_name') or path.name}`",
        f"- 结构来源：`{plan.get('structure_source')}`",
        f"- 置信度：`{plan.get('confidence')}`",
        f"- 默认抽取：`{plan.get('default_extraction', 'all_items')}`",
        f"- 内容单元数：`{len(items)}`",
        "",
        "| 序号 | ID | 标题 | 范围 | tokens | 状态 | Markdown |",
        "| ---: | --- | --- | --- | ---: | --- | --- |",
    ]

    for index, item in enumerate(items, start=1):
        item_id = str(item.get("id") or f"item{index:03d}")
        title = str(item.get("title") or item_id)
        scope = item_scope(item)
        filename = f"{index:03d}-{slugify_filename(item_id, f'item{index:03d}', 24)}-{slugify_filename(title, 'untitled', 48)}.md"
        item_path = out_dir / filename
        body = extract_item(path, item, max_chars).strip()
        content = "\n".join(
            [
                f"# {title}",
                "",
                f"> 来源文件：{plan.get('file_name') or path.name}",
                f"> 内容单元：{item_id}；范围：{scope}",
                f"> 结构来源：{plan.get('structure_source')}；置信度：{plan.get('confidence')}",
                f"> 估算 tokens：{item.get('estimated_tokens', 0)}；状态：{item.get('status', '')}",
                "",
                body,
                "",
            ]
        )
        item_path.write_text(content, encoding="utf-8")
        index_lines.append(
            f"| {index} | `{markdown_cell(item_id)}` | {markdown_cell(title)} | {markdown_cell(scope)} | "
            f"{item.get('estimated_tokens', 0)} | {markdown_cell(item.get('status'))} | [{markdown_cell(filename)}]({filename}) |"
        )

    if not items:
        index_lines.append("|  |  | 未找到可抽取内容单元 |  |  |  |  |")

    warnings = plan.get("warnings") or []
    index_lines.extend(["", "## 规划警告", ""])
    index_lines.extend([f"- {warning}" for warning in warnings] or ["- 无"])
    index_lines.append("")
    index_path = out_dir / "index.md"
    index_text = "\n".join(index_lines)
    index_path.write_text(index_text, encoding="utf-8")
    return index_text + f"\n<!-- output-dir: {out_dir} -->\n"


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Extract selected document ranges with source anchors.")
    parser.add_argument("file")
    parser.add_argument("--pages", help="PDF page range, e.g. 1-3,8.")
    parser.add_argument("--slides", help="PPTX slide range, e.g. 2-5.")
    parser.add_argument("--paragraphs", help="DOCX/Markdown/TXT paragraph range, e.g. 1-10.")
    parser.add_argument("--heading", help="Extract a DOCX/Markdown/TXT section by heading text.")
    parser.add_argument("--plan-json", help="Chapter/chunk plan JSON from plan_chapter_batches.py.")
    parser.add_argument("--batch-id", help="Batch ID to extract from --plan-json, or 'first'.")
    parser.add_argument("--all-items", action="store_true", help="Extract every planned item from --plan-json.")
    parser.add_argument("--output-dir", help="Directory for one Markdown file per planned item; use with --plan-json --all-items.")
    parser.add_argument("--max-chars", type=int, default=200000, help="Maximum characters to output.")
    parser.add_argument("--output", help="Optional Markdown output path.")
    args = parser.parse_args(argv)

    if args.all_items and not args.plan_json:
        raise SystemExit("--all-items requires --plan-json")
    if args.output_dir and not (args.plan_json and args.all_items):
        raise SystemExit("--output-dir requires --plan-json --all-items")

    path = Path(args.file)
    if args.plan_json and args.all_items and args.output_dir:
        output = extract_plan_items_to_dir(path, args.plan_json, args.output_dir, args.max_chars)
    else:
        output = extract_range(path, args)
    if args.output:
        Path(args.output).write_text(output, encoding="utf-8")
    else:
        sys.stdout.write(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
