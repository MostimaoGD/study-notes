#!/usr/bin/env python
"""Audit learning documents before asking an LLM to study them."""

from __future__ import annotations

import argparse
import json
import math
import re
import sys
from pathlib import Path
from typing import Any


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".md", ".markdown", ".txt"}
CJK_RE = re.compile(r"[\u3400-\u9fff\uf900-\ufaff]")
WORD_RE = re.compile(r"[A-Za-z0-9_]+(?:[-'][A-Za-z0-9_]+)?")


def estimate_tokens(text: str) -> int:
    """Return a conservative token estimate for mixed Chinese/English text."""
    if not text:
        return 0
    cjk_count = len(CJK_RE.findall(text))
    non_cjk_text = CJK_RE.sub(" ", text)
    word_count = len(WORD_RE.findall(non_cjk_text))
    punctuation_count = sum(1 for ch in text if not ch.isspace() and not ch.isalnum())
    char_estimate = math.ceil(len(text) / 3.6)
    mixed_estimate = math.ceil(cjk_count * 1.15 + word_count * 1.35 + punctuation_count * 0.2)
    return max(char_estimate, mixed_estimate)


def read_text_file(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "utf-16"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def text_quality_for_pdf(total_chars: int, pages: int, low_text_pages: int) -> str:
    if pages <= 0 or total_chars < 50:
        return "poor"
    avg_chars = total_chars / pages
    low_ratio = low_text_pages / pages
    if avg_chars >= 800 and low_ratio <= 0.25:
        return "good"
    if avg_chars >= 200 and low_ratio <= 0.60:
        return "fair"
    return "poor"


def text_quality_for_plain(total_chars: int) -> str:
    if total_chars >= 1200:
        return "good"
    if total_chars >= 120:
        return "fair"
    return "poor"


def fit_status(text_quality: str, scan_risk: str, tokens: int, context_limit: int, reserve: int) -> str:
    available = max(1000, context_limit - reserve)
    if tokens == 0:
        return "unreadable"
    if text_quality == "poor" and scan_risk in {"medium", "high"}:
        return "unreadable"
    if tokens > available:
        return "oversize"
    if tokens > available * 0.8 or text_quality in {"fair", "poor"} or scan_risk == "medium":
        return "borderline"
    return "fits"


def audit_pdf(path: Path) -> dict[str, Any]:
    result: dict[str, Any] = {
        "file_type": "pdf",
        "pages": 0,
        "page_basis": "actual",
        "text": "",
        "total_chars": 0,
        "text_pages": 0,
        "low_text_pages": 0,
        "image_count": 0,
        "image_pages": 0,
        "table_count": 0,
        "warnings": [],
        "errors": [],
    }

    try:
        import pdfplumber
    except ImportError as exc:
        result["errors"].append(f"pdfplumber unavailable: {exc}")
        return result

    try:
        with pdfplumber.open(str(path)) as pdf:
            result["pages"] = len(pdf.pages)
            page_texts: list[str] = []
            for index, page in enumerate(pdf.pages, start=1):
                try:
                    text = page.extract_text() or ""
                except Exception as exc:  # pragma: no cover - depends on malformed PDFs
                    text = ""
                    result["warnings"].append(f"p.{index}: text extraction failed: {exc}")

                chars = len(text.strip())
                if chars:
                    result["text_pages"] += 1
                if chars < 80:
                    result["low_text_pages"] += 1
                page_texts.append(f"<!-- source: {path.name} p.{index} -->\n{text}".strip())

                images = getattr(page, "images", None) or []
                result["image_count"] += len(images)
                if images:
                    result["image_pages"] += 1

                try:
                    tables = page.extract_tables() or []
                    result["table_count"] += len(tables)
                except Exception:
                    pass

            result["text"] = "\n\n".join(page_texts)
            result["total_chars"] = len(result["text"])
    except Exception as exc:
        result["errors"].append(f"PDF open/extract failed: {exc}")
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            result["pages"] = len(reader.pages)
            fallback_texts = []
            for index, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ""
                fallback_texts.append(f"<!-- source: {path.name} p.{index} -->\n{text}".strip())
                if text.strip():
                    result["text_pages"] += 1
                if len(text.strip()) < 80:
                    result["low_text_pages"] += 1
            result["text"] = "\n\n".join(fallback_texts)
            result["total_chars"] = len(result["text"])
            result["warnings"].append("Used pypdf fallback extraction.")
        except Exception as fallback_exc:
            result["errors"].append(f"pypdf fallback failed: {fallback_exc}")

    pages = int(result["pages"] or 0)
    total_chars = int(result["total_chars"] or 0)
    image_pages = int(result["image_pages"] or 0)
    low_text_pages = int(result["low_text_pages"] or 0)
    image_ratio = image_pages / pages if pages else 0
    low_ratio = low_text_pages / pages if pages else 1
    avg_chars = total_chars / pages if pages else 0

    result["text_quality"] = text_quality_for_pdf(total_chars, pages, low_text_pages)
    if image_ratio >= 0.5 and avg_chars < 200:
        result["scan_risk"] = "high"
    elif image_ratio >= 0.3 or (image_pages > 0 and low_ratio >= 0.5):
        result["scan_risk"] = "medium"
    else:
        result["scan_risk"] = "low"
    return result


def audit_docx(path: Path) -> dict[str, Any]:
    result: dict[str, Any] = {
        "file_type": "docx",
        "pages": 0,
        "page_basis": "estimated",
        "text": "",
        "total_chars": 0,
        "paragraphs": 0,
        "tables": 0,
        "image_count": 0,
        "warnings": [],
        "errors": [],
    }
    try:
        from docx import Document

        doc = Document(str(path))
        parts: list[str] = []
        for index, paragraph in enumerate(doc.paragraphs, start=1):
            text = paragraph.text.strip()
            if text:
                parts.append(f"<!-- source: {path.name} paragraph {index} -->\n{text}")
        for table_index, table in enumerate(doc.tables, start=1):
            rows = []
            for row in table.rows:
                rows.append(" | ".join(cell.text.strip() for cell in row.cells))
            if rows:
                parts.append(f"<!-- source: {path.name} table {table_index} -->\n" + "\n".join(rows))
        result["paragraphs"] = len(doc.paragraphs)
        result["tables"] = len(doc.tables)
        result["image_count"] = len(doc.inline_shapes)
        result["text"] = "\n\n".join(parts)
        result["total_chars"] = len(result["text"])
        result["pages"] = max(1, math.ceil(result["total_chars"] / 1800)) if result["total_chars"] else 0
    except Exception as exc:
        result["errors"].append(f"DOCX extraction failed: {exc}")

    result["text_quality"] = text_quality_for_plain(int(result["total_chars"] or 0))
    result["scan_risk"] = "low"
    return result


def iter_shape_text(shape: Any) -> list[str]:
    texts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip()
        if text:
            texts.append(text)
    if getattr(shape, "has_table", False):
        try:
            rows = []
            for row in shape.table.rows:
                rows.append(" | ".join(cell.text.strip() for cell in row.cells))
            if rows:
                texts.append("\n".join(rows))
        except Exception:
            pass
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            texts.extend(iter_shape_text(child))
    return texts


def audit_pptx(path: Path) -> dict[str, Any]:
    result: dict[str, Any] = {
        "file_type": "pptx",
        "pages": 0,
        "page_basis": "slides",
        "text": "",
        "total_chars": 0,
        "slides": 0,
        "image_count": 0,
        "table_count": 0,
        "chart_count": 0,
        "warnings": [],
        "errors": [],
    }
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(str(path))
        result["slides"] = len(prs.slides)
        result["pages"] = len(prs.slides)
        slide_texts: list[str] = []
        for slide_index, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                texts.extend(iter_shape_text(shape))
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    result["image_count"] += 1
                if getattr(shape, "has_table", False):
                    result["table_count"] += 1
                if getattr(shape, "has_chart", False):
                    result["chart_count"] += 1
            if texts:
                slide_texts.append(f"<!-- source: {path.name} slide {slide_index} -->\n" + "\n\n".join(texts))
        result["text"] = "\n\n".join(slide_texts)
        result["total_chars"] = len(result["text"])
    except Exception as exc:
        result["errors"].append(f"PPTX extraction failed: {exc}")

    result["text_quality"] = text_quality_for_plain(int(result["total_chars"] or 0))
    result["scan_risk"] = "medium" if result.get("image_count", 0) and result["text_quality"] == "poor" else "low"
    return result


def audit_plain(path: Path, file_type: str) -> dict[str, Any]:
    result: dict[str, Any] = {
        "file_type": file_type,
        "pages": 0,
        "page_basis": "estimated",
        "text": "",
        "total_chars": 0,
        "warnings": [],
        "errors": [],
    }
    try:
        text = read_text_file(path)
        result["text"] = f"<!-- source: {path.name} -->\n{text}"
        result["total_chars"] = len(text)
        result["pages"] = max(1, math.ceil(len(text) / 1800)) if text else 0
    except Exception as exc:
        result["errors"].append(f"Plain text extraction failed: {exc}")
    result["text_quality"] = text_quality_for_plain(int(result["total_chars"] or 0))
    result["scan_risk"] = "low"
    return result


def audit_document(path: Path, context_limit: int, reserve: int) -> dict[str, Any]:
    ext = path.suffix.lower()
    base: dict[str, Any] = {
        "file": str(path),
        "file_name": path.name,
        "file_type": ext.lstrip(".") or "unknown",
        "pages": 0,
        "estimated_tokens": 0,
        "text_quality": "poor",
        "scan_risk": "low",
        "fit_status": "unreadable",
        "context_limit": context_limit,
        "reserved_tokens": reserve,
        "available_tokens": max(1000, context_limit - reserve),
        "warnings": [],
        "errors": [],
    }

    if not path.exists():
        base["errors"].append("File does not exist.")
        return base
    if ext not in SUPPORTED_EXTENSIONS:
        base["errors"].append(f"Unsupported extension: {ext}")
        return base

    if ext == ".pdf":
        detail = audit_pdf(path)
    elif ext == ".docx":
        detail = audit_docx(path)
    elif ext == ".pptx":
        detail = audit_pptx(path)
    elif ext in {".md", ".markdown"}:
        detail = audit_plain(path, "markdown")
    else:
        detail = audit_plain(path, "txt")

    text = detail.pop("text", "")
    detail_warnings = list(detail.pop("warnings", []))
    detail_errors = list(detail.pop("errors", []))
    base.update(detail)
    base["warnings"] = list(base.get("warnings", [])) + detail_warnings
    base["errors"] = list(base.get("errors", [])) + detail_errors
    base["estimated_tokens"] = estimate_tokens(text)
    base["total_chars"] = len(text)
    base["fit_status"] = fit_status(
        str(base.get("text_quality", "poor")),
        str(base.get("scan_risk", "low")),
        int(base["estimated_tokens"]),
        context_limit,
        reserve,
    )
    if base["fit_status"] == "unreadable":
        base["recommendation"] = "不要声称已完整读取；请获取可复制文本版本、OCR 后文本，或让用户指定可读范围。"
    elif base["fit_status"] == "oversize":
        base["recommendation"] = "文件过长；请先运行章节/分块规划并抽取全部内容单元，再按章节/分块生成 Markdown 笔记。不要尝试一次性把全文塞入上下文。"
    elif base["fit_status"] == "borderline":
        base["recommendation"] = "请先运行章节/分块规划并抽取全部内容单元，再按章节/分块生成 Markdown 笔记；必要时说明抽取质量风险。"
    else:
        base["recommendation"] = "可以继续运行章节/分块规划并抽取全部内容单元，然后生成有证据链的 Markdown 图文笔记。"
    return base


def render_markdown(result: dict[str, Any]) -> str:
    warning_lines = "\n".join(f"- {item}" for item in result.get("warnings", []) if item) or "- 无"
    error_lines = "\n".join(f"- {item}" for item in result.get("errors", []) if item) or "- 无"
    return "\n".join(
        [
            "# 文档可读性审计",
            "",
            f"- 文件：`{result.get('file_name', '')}`",
            f"- 类型：`{result.get('file_type', '')}`",
            f"- 页数/页等价：`{result.get('pages', 0)}`（{result.get('page_basis', 'unknown')}）",
            f"- 字符数：`{result.get('total_chars', 0)}`",
            f"- 估算 tokens：`{result.get('estimated_tokens', 0)}`",
            f"- 文本抽取质量：`{result.get('text_quality', '')}`",
            f"- 扫描/图片风险：`{result.get('scan_risk', '')}`",
            f"- 上下文预算：`{result.get('available_tokens', 0)}` 可用 / `{result.get('context_limit', 0)}` 总量",
            f"- 读取结论：`{result.get('fit_status', '')}`",
            "",
            "## 建议",
            "",
            str(result.get("recommendation", "")),
            "",
            "## 警告",
            "",
            warning_lines,
            "",
            "## 错误",
            "",
            error_lines,
            "",
        ]
    )


def write_text(path: str | None, content: str) -> None:
    if path:
        Path(path).write_text(content, encoding="utf-8")


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Audit whether a learning document can be read safely.")
    parser.add_argument("file", help="PDF, DOCX, PPTX, Markdown, or TXT file to audit.")
    parser.add_argument("--context-limit", type=int, default=100000, help="Total model context budget.")
    parser.add_argument("--reserve", type=int, default=30000, help="Tokens reserved for prompts, reasoning, and output.")
    parser.add_argument("--json-output", help="Optional path for JSON audit output.")
    parser.add_argument("--markdown-output", help="Optional path for Markdown audit output.")
    parser.add_argument("--print-json", action="store_true", help="Print JSON instead of Markdown to stdout.")
    args = parser.parse_args(argv)

    result = audit_document(Path(args.file), args.context_limit, args.reserve)
    json_text = json.dumps(result, ensure_ascii=False, indent=2)
    markdown_text = render_markdown(result)
    write_text(args.json_output, json_text + "\n")
    write_text(args.markdown_output, markdown_text)
    print(json_text if args.print_json else markdown_text)
    return 0 if not result.get("errors") else 1


if __name__ == "__main__":
    raise SystemExit(main())
