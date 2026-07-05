#!/usr/bin/env python
"""Plan chapters, sections, or chunks for learning documents."""

from __future__ import annotations

import argparse
import json
import math
import re
import statistics
from pathlib import Path
from typing import Any, Iterable


CJK_RE = re.compile(r"[\u3400-\u9fff\uf900-\ufaff]")
WORD_RE = re.compile(r"[A-Za-z0-9_]+(?:[-'][A-Za-z0-9_]+)?")
CHAPTER_TITLE_RE = re.compile(
    r"^\s*(第\s*[0-9一二三四五六七八九十百千万]+\s*章\s+[^\n·.]{1,80}|Chapter\s+\d+[:.\-\s]+[^\n]{1,80})\s*$",
    re.IGNORECASE,
)
CHINESE_TOC_RE = re.compile(
    r"^\s*(第\s*[0-9一二三四五六七八九十百千万]+\s*章\s+[^\n·.]{1,80})\s*[·.\-\s]{2,}[\(（]?\s*(\d{1,4})\s*[\)）]?\s*$"
)
APPENDIX_TOC_RE = re.compile(
    r"^\s*(附录\s*[A-ZＡ-Ｚ]?\s*[^\n·.]{0,80})\s*[·.\-\s]{2,}[\(（]?\s*(\d{1,4})\s*[\)）]?\s*$",
    re.IGNORECASE,
)
ENGLISH_TOC_RE = re.compile(
    r"^\s*((?:Chapter|CHAPTER)\s+\d+[:.\-\s]+[^\n]{1,80}?)\s*[.\-\s]{2,}(\d{1,4})\s*$"
)
NUMBERED_HEADING_RE = re.compile(r"^\s*(\d+(?:\.\d+){0,2})[.\s]+([A-Z][^\n]{2,80}|[\u4e00-\u9fff][^\n]{2,80})\s*$")
PAPER_SECTION_RE = re.compile(
    r"^\s*(?:\d+\.?\s+)?("
    r"Abstract|Introduction|Background|Related\s+Work|Methods?|Materials\s+and\s+Methods|Methodology|"
    r"Experiments?|Results?|Discussion|Conclusion|Conclusions|References|Acknowledg(?:e)?ments?|"
    r"摘要|引言|绪论|背景|相关工作|方法|实验|结果|讨论|结论|参考文献|致谢"
    r")\s*$",
    re.IGNORECASE,
)


def estimate_tokens(text: str) -> int:
    if not text:
        return 0
    cjk_count = len(CJK_RE.findall(text))
    non_cjk_text = CJK_RE.sub(" ", text)
    word_count = len(WORD_RE.findall(non_cjk_text))
    punctuation_count = sum(1 for ch in text if not ch.isspace() and not ch.isalnum())
    char_estimate = math.ceil(len(text) / 3.6)
    mixed_estimate = math.ceil(cjk_count * 1.15 + word_count * 1.35 + punctuation_count * 0.2)
    return max(char_estimate, mixed_estimate)


def read_text_file(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "utf-16"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def split_plain_blocks(text: str, max_block_chars: int = 6000) -> list[str]:
    blocks: list[str] = []
    for part in re.split(r"\n\s*\n", text):
        part = part.strip()
        if not part:
            continue
        if len(part) <= max_block_chars:
            blocks.append(part)
            continue
        lines = [line.strip() for line in part.splitlines() if line.strip()]
        if len(lines) > 1 and max(len(line) for line in lines) <= max_block_chars:
            blocks.extend(lines)
            continue
        for start in range(0, len(part), max_block_chars):
            blocks.append(part[start : start + max_block_chars].strip())
    return [block for block in blocks if block]


def normalize(text: str) -> str:
    return re.sub(r"[^0-9A-Za-z\u4e00-\u9fff]+", "", text).lower()


def range_spec(start: int, end: int) -> str:
    return str(start) if start == end else f"{start}-{end}"


def compact_title(text: str, limit: int = 80) -> str:
    text = re.sub(r"\s+", " ", text).strip(" .·\t")
    return text[:limit].strip()


def iter_lines_with_page(page_texts: list[str]) -> Iterable[tuple[int, str]]:
    for page_number, text in enumerate(page_texts, start=1):
        for raw_line in text.splitlines():
            line = compact_title(raw_line, 140)
            if line:
                yield page_number, line


def extract_pdf_pages(path: Path) -> tuple[list[str], list[int], list[str]]:
    warnings: list[str] = []
    try:
        import pdfplumber
    except ImportError as exc:
        raise SystemExit(f"pdfplumber unavailable: {exc}")

    page_texts: list[str] = []
    page_tokens: list[int] = []
    with pdfplumber.open(str(path)) as pdf:
        for page_number, page in enumerate(pdf.pages, start=1):
            try:
                text = page.extract_text() or ""
            except Exception as exc:  # pragma: no cover - depends on malformed PDFs
                text = ""
                warnings.append(f"p.{page_number}: text extraction failed: {exc}")
            page_texts.append(text)
            page_tokens.append(estimate_tokens(text))
    return page_texts, page_tokens, warnings


def flatten_outline(outline: Any, reader: Any, level: int = 0) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    for entry in outline or []:
        if isinstance(entry, list):
            items.extend(flatten_outline(entry, reader, level + 1))
            continue
        title = getattr(entry, "title", None)
        if not title:
            continue
        try:
            page = reader.get_destination_page_number(entry) + 1
        except Exception:
            page = None
        if page:
            items.append({"title": compact_title(str(title)), "start_page": page, "level": level})
    return items


def read_pdf_bookmarks(path: Path, total_pages: int) -> list[dict[str, Any]]:
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        outline = getattr(reader, "outline", None) or getattr(reader, "outlines", None)
        raw = flatten_outline(outline, reader)
    except Exception:
        return []

    chapterish = [
        item
        for item in raw
        if item["level"] <= 2 and (CHAPTER_TITLE_RE.match(item["title"]) or PAPER_SECTION_RE.match(item["title"]))
    ]
    if len(chapterish) < 2:
        top_level = [item for item in raw if item["level"] == 0]
        chapterish = top_level if len(top_level) >= 2 else []
    chapterish = sorted(chapterish, key=lambda item: item["start_page"])
    return finish_page_items(chapterish, total_pages, "bookmark")


def parse_visible_toc(page_texts: list[str], total_pages: int) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    scan_pages = min(total_pages, 20)
    for page_index in range(scan_pages):
        for line in page_texts[page_index].splitlines():
            line = compact_title(line, 180)
            if not line:
                continue
            match = CHINESE_TOC_RE.match(line) or APPENDIX_TOC_RE.match(line) or ENGLISH_TOC_RE.match(line)
            if not match:
                continue
            title = compact_title(match.group(1))
            declared_page = int(match.group(2))
            if any(normalize(item["title"]) == normalize(title) for item in items):
                continue
            items.append({"title": title, "declared_page": declared_page, "toc_pdf_page": page_index + 1})
    return items


def find_title_page(title: str, page_texts: list[str], min_page: int = 1) -> int | None:
    title_norm = normalize(title)
    if not title_norm:
        return None
    title_tail = title_norm
    chapter_prefix = re.match(r"(第[0-9一二三四五六七八九十百千万]+章)(.+)", title_norm)
    if chapter_prefix:
        title_tail = chapter_prefix.group(2)
    for page_number, text in enumerate(page_texts, start=1):
        if page_number < min_page:
            continue
        page_norm = normalize(text[:2500])
        if title_norm and title_norm in page_norm:
            return page_number
        if chapter_prefix and title_tail and len(title_tail) >= 3:
            first_lines = "\n".join(text.splitlines()[:12])
            for raw_line in first_lines.splitlines():
                line_norm = normalize(raw_line)
                if len(line_norm) <= 80 and title_tail in line_norm and chapter_prefix.group(1) in line_norm:
                    return page_number
    return None


def apply_toc_page_offset(toc_items: list[dict[str, Any]], page_texts: list[str], total_pages: int) -> tuple[list[dict[str, Any]], int | None, list[str]]:
    warnings: list[str] = []
    offsets: list[int] = []
    search_min_page = max(int(item.get("toc_pdf_page") or 1) for item in toc_items) + 1
    for item in toc_items:
        actual = find_title_page(item["title"], page_texts, search_min_page)
        if actual:
            item["start_page"] = actual
            offsets.append(actual - int(item["declared_page"]))
    page_offset = int(statistics.median(offsets)) if offsets else None
    if page_offset is not None:
        for item in toc_items:
            item.setdefault("start_page", int(item["declared_page"]) + page_offset)
    else:
        warnings.append("Could not infer a reliable PDF page offset from the visible table of contents.")

    valid = []
    for item in toc_items:
        start = item.get("start_page")
        if isinstance(start, int) and 1 <= start <= total_pages:
            valid.append(item)
        else:
            warnings.append(f"Skipped TOC item with unresolved page: {item.get('title')}")
    return finish_page_items(valid, total_pages, "toc"), page_offset, warnings


def detect_pdf_headings(page_texts: list[str], total_pages: int, paper_only: bool = False) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    seen: set[str] = set()
    for page_number, line in iter_lines_with_page(page_texts):
        if "··" in line or "..." in line or len(line) > 120:
            continue
        candidates: list[str] = []
        if not paper_only:
            if CHAPTER_TITLE_RE.match(line):
                candidates.append(line)
            elif NUMBERED_HEADING_RE.match(line) and page_number > 1:
                number = NUMBERED_HEADING_RE.match(line).group(1)
                if number.count(".") == 0:
                    candidates.append(line)
        if paper_only and PAPER_SECTION_RE.match(line):
            candidates.append(line)
        for title in candidates:
            key = normalize(title)
            if key and key not in seen:
                items.append({"title": title, "start_page": page_number})
                seen.add(key)
    return finish_page_items(sorted(items, key=lambda item: item["start_page"]), total_pages, "heading")


def finish_page_items(items: list[dict[str, Any]], total_pages: int, source: str) -> list[dict[str, Any]]:
    cleaned: list[dict[str, Any]] = []
    last_start = 0
    for item in sorted(items, key=lambda value: value.get("start_page") or 0):
        start = int(item.get("start_page") or 0)
        if start <= last_start or start < 1 or start > total_pages:
            continue
        new_item = dict(item)
        new_item["start_page"] = start
        new_item["source"] = source
        cleaned.append(new_item)
        last_start = start
    for index, item in enumerate(cleaned):
        next_start = cleaned[index + 1]["start_page"] if index + 1 < len(cleaned) else total_pages + 1
        item["end_page"] = max(item["start_page"], next_start - 1)
    return cleaned


def page_tokens_for_range(page_tokens: list[int], start: int, end: int) -> int:
    return sum(page_tokens[start - 1 : end])


def make_pdf_item(raw: dict[str, Any], index: int, page_tokens: list[int], budget: int) -> dict[str, Any]:
    start = int(raw["start_page"])
    end = int(raw["end_page"])
    tokens = page_tokens_for_range(page_tokens, start, end)
    status = "chapter_oversize" if tokens > budget else "fits"
    kind = "paper_section" if PAPER_SECTION_RE.match(raw.get("title", "")) else "chapter"
    if raw.get("title", "").lower().startswith("附录"):
        kind = "appendix"
    return {
        "id": f"ch{index:03d}",
        "title": raw.get("title") or f"Pages {start}-{end}",
        "kind": kind,
        "start_page": start,
        "end_page": end,
        "estimated_tokens": tokens,
        "status": status,
        "extraction": {"pages": range_spec(start, end)},
    }


def split_pdf_item(item: dict[str, Any], page_tokens: list[int], budget: int) -> list[dict[str, Any]]:
    if item["estimated_tokens"] <= budget:
        return [item]
    pieces: list[dict[str, Any]] = []
    start = int(item["start_page"])
    end = int(item["end_page"])
    current_start = start
    current_tokens = 0
    part = 1
    for page in range(start, end + 1):
        token_count = page_tokens[page - 1]
        if current_tokens and current_tokens + token_count > budget:
            current_end = page - 1
            pieces.append(
                make_split_item(item, part, current_start, current_end, current_tokens, "pages", range_spec(current_start, current_end))
            )
            part += 1
            current_start = page
            current_tokens = 0
        current_tokens += token_count
    if current_start <= end:
        pieces.append(make_split_item(item, part, current_start, end, current_tokens, "pages", range_spec(current_start, end)))
    return pieces


def make_split_item(parent: dict[str, Any], part: int, start: int, end: int, tokens: int, extraction_key: str, extraction_value: str) -> dict[str, Any]:
    return {
        "id": f"{parent['id']}-part{part}",
        "title": f"{parent['title']}（part {part}）",
        "kind": "chapter_chunk",
        "parent_id": parent["id"],
        "parent_title": parent["title"],
        "start_page": start if extraction_key == "pages" else None,
        "end_page": end if extraction_key == "pages" else None,
        "start_paragraph": start if extraction_key == "paragraphs" else None,
        "end_paragraph": end if extraction_key == "paragraphs" else None,
        "estimated_tokens": tokens,
        "status": "chapter_oversize",
        "extraction": {extraction_key: extraction_value},
    }


def make_auto_page_chunks(page_tokens: list[int], total_pages: int, budget: int) -> list[dict[str, Any]]:
    chunks: list[dict[str, Any]] = []
    start = 1
    current = 0
    index = 1
    for page in range(1, total_pages + 1):
        token_count = page_tokens[page - 1]
        if current and current + token_count > budget:
            end = page - 1
            chunks.append(
                {
                    "id": f"chunk{index:03d}",
                    "title": f"自动分块 {index}: p.{start}-{end}",
                    "kind": "auto_chunk",
                    "start_page": start,
                    "end_page": end,
                    "estimated_tokens": current,
                    "status": "fits",
                    "extraction": {"pages": range_spec(start, end)},
                }
            )
            index += 1
            start = page
            current = 0
        current += token_count
    if start <= total_pages:
        chunks.append(
            {
                "id": f"chunk{index:03d}",
                "title": f"自动分块 {index}: p.{start}-{total_pages}",
                "kind": "auto_chunk",
                "start_page": start,
                "end_page": total_pages,
                "estimated_tokens": current,
                "status": "fits" if current <= budget else "chunk_oversize",
                "extraction": {"pages": range_spec(start, total_pages)},
            }
        )
    return chunks


def plan_pdf(path: Path, context_limit: int, reserve: int, item_budget: int) -> dict[str, Any]:
    page_texts, page_tokens, warnings = extract_pdf_pages(path)
    total_pages = len(page_texts)
    page_offset: int | None = None

    raw = read_pdf_bookmarks(path, total_pages)
    structure_source = "bookmarks"
    confidence = "high"
    if raw:
        toc_items = parse_visible_toc(page_texts, total_pages)
        if len(toc_items) >= 2:
            _, page_offset, _ = apply_toc_page_offset(toc_items, page_texts, total_pages)
    if not raw:
        toc_items = parse_visible_toc(page_texts, total_pages)
        if len(toc_items) >= 2:
            raw, page_offset, toc_warnings = apply_toc_page_offset(toc_items, page_texts, total_pages)
            warnings.extend(toc_warnings)
            structure_source = "toc"
            confidence = "high" if raw and page_offset is not None else "medium"
    if not raw:
        raw = detect_pdf_headings(page_texts, total_pages)
        structure_source = "headings"
        confidence = "medium" if len(raw) >= 2 else "low"
    if len(raw) < 2:
        raw = detect_pdf_headings(page_texts, total_pages, paper_only=True)
        structure_source = "paper_sections"
        confidence = "medium" if len(raw) >= 2 else "low"

    if len(raw) >= 2:
        base_items = [make_pdf_item(item, index + 1, page_tokens, item_budget) for index, item in enumerate(raw)]
        items: list[dict[str, Any]] = []
        for item in base_items:
            if item["status"] == "chapter_oversize":
                warnings.append(f"Chapter exceeds batch budget and was split: {item['title']}")
                items.extend(split_pdf_item(item, page_tokens, item_budget))
            else:
                items.append(item)
    else:
        items = make_auto_page_chunks(page_tokens, total_pages, item_budget)
        structure_source = "auto_chunks"
        confidence = "low"
        warnings.append("No reliable table of contents or headings found; planned automatic page/token chunks.")

    return make_plan(path, "pdf", structure_source, confidence, page_offset, items, context_limit, reserve, warnings)


def docx_paragraphs(path: Path) -> list[dict[str, Any]]:
    from docx import Document

    doc = Document(str(path))
    paragraphs = []
    for index, paragraph in enumerate(doc.paragraphs, start=1):
        text = paragraph.text.strip()
        style = paragraph.style.name if paragraph.style else ""
        if text:
            paragraphs.append({"number": index, "text": text, "style": style})
    return paragraphs


def detect_paragraph_items(paragraphs: list[dict[str, Any]], total_paragraphs: int, source_type: str) -> tuple[list[dict[str, Any]], str, str]:
    headings: list[dict[str, Any]] = []
    for paragraph in paragraphs:
        text = paragraph["text"]
        style = str(paragraph.get("style", "")).lower()
        is_heading = False
        if source_type == "docx" and re.search(r"heading\s+1|标题\s*1", style):
            is_heading = True
        elif source_type in {"markdown", "txt"} and text.startswith("# "):
            is_heading = True
            text = text.lstrip("# ").strip()
        elif CHAPTER_TITLE_RE.match(text) or PAPER_SECTION_RE.match(text):
            is_heading = True
        if is_heading:
            headings.append({"title": compact_title(text), "start_paragraph": int(paragraph["number"])})

    if len(headings) < 2:
        paper_sections = []
        for paragraph in paragraphs:
            text = paragraph["text"]
            if PAPER_SECTION_RE.match(text):
                paper_sections.append({"title": compact_title(text), "start_paragraph": int(paragraph["number"])})
        headings = paper_sections
        structure_source = "paper_sections"
        confidence = "medium" if len(headings) >= 2 else "low"
    else:
        structure_source = "headings"
        confidence = "medium"

    cleaned = []
    for index, item in enumerate(headings):
        start = item["start_paragraph"]
        end = (headings[index + 1]["start_paragraph"] - 1) if index + 1 < len(headings) else total_paragraphs
        cleaned.append({**item, "end_paragraph": end, "source": structure_source})
    return cleaned, structure_source, confidence


def text_to_paragraphs(path: Path) -> list[dict[str, Any]]:
    text = read_text_file(path)
    raw_parts = split_plain_blocks(text)
    paragraphs: list[dict[str, Any]] = []
    for index, part in enumerate(raw_parts, start=1):
        first_line = part.splitlines()[0].strip()
        paragraphs.append({"number": index, "text": first_line, "full_text": part, "style": ""})
    return paragraphs


def paragraph_tokens(paragraphs: list[dict[str, Any]], start: int, end: int) -> int:
    token_total = 0
    by_number = {int(item["number"]): item for item in paragraphs}
    for number in range(start, end + 1):
        item = by_number.get(number)
        if item:
            token_total += estimate_tokens(item.get("full_text") or item.get("text") or "")
    return token_total


def make_paragraph_items(raw: list[dict[str, Any]], paragraphs: list[dict[str, Any]], budget: int) -> list[dict[str, Any]]:
    items = []
    for index, item in enumerate(raw, start=1):
        start = int(item["start_paragraph"])
        end = int(item["end_paragraph"])
        tokens = paragraph_tokens(paragraphs, start, end)
        status = "chapter_oversize" if tokens > budget else "fits"
        kind = "paper_section" if PAPER_SECTION_RE.match(item["title"]) else "chapter"
        base = {
            "id": f"ch{index:03d}",
            "title": item["title"],
            "kind": kind,
            "start_paragraph": start,
            "end_paragraph": end,
            "estimated_tokens": tokens,
            "status": status,
            "extraction": {"paragraphs": range_spec(start, end)},
        }
        if status == "chapter_oversize":
            items.extend(split_paragraph_item(base, paragraphs, budget))
        else:
            items.append(base)
    return items


def split_paragraph_item(item: dict[str, Any], paragraphs: list[dict[str, Any]], budget: int) -> list[dict[str, Any]]:
    pieces: list[dict[str, Any]] = []
    by_number = {int(value["number"]): value for value in paragraphs}
    current_start = int(item["start_paragraph"])
    current_tokens = 0
    part = 1
    for number in range(int(item["start_paragraph"]), int(item["end_paragraph"]) + 1):
        token_count = estimate_tokens((by_number.get(number) or {}).get("full_text") or (by_number.get(number) or {}).get("text") or "")
        if current_tokens and current_tokens + token_count > budget:
            pieces.append(make_split_item(item, part, current_start, number - 1, current_tokens, "paragraphs", range_spec(current_start, number - 1)))
            part += 1
            current_start = number
            current_tokens = 0
        current_tokens += token_count
    pieces.append(make_split_item(item, part, current_start, int(item["end_paragraph"]), current_tokens, "paragraphs", range_spec(current_start, int(item["end_paragraph"]))))
    return pieces


def make_auto_paragraph_chunks(paragraphs: list[dict[str, Any]], budget: int) -> list[dict[str, Any]]:
    chunks: list[dict[str, Any]] = []
    current_start = int(paragraphs[0]["number"]) if paragraphs else 1
    current_tokens = 0
    index = 1
    for paragraph in paragraphs:
        number = int(paragraph["number"])
        token_count = estimate_tokens(paragraph.get("full_text") or paragraph.get("text") or "")
        if current_tokens and current_tokens + token_count > budget:
            end = number - 1
            chunks.append(
                {
                    "id": f"chunk{index:03d}",
                    "title": f"自动分块 {index}: paragraph {current_start}-{end}",
                    "kind": "auto_chunk",
                    "start_paragraph": current_start,
                    "end_paragraph": end,
                    "estimated_tokens": current_tokens,
                    "status": "fits",
                    "extraction": {"paragraphs": range_spec(current_start, end)},
                }
            )
            index += 1
            current_start = number
            current_tokens = 0
        current_tokens += token_count
    if paragraphs:
        end = int(paragraphs[-1]["number"])
        chunks.append(
            {
                "id": f"chunk{index:03d}",
                "title": f"自动分块 {index}: paragraph {current_start}-{end}",
                "kind": "auto_chunk",
                "start_paragraph": current_start,
                "end_paragraph": end,
                "estimated_tokens": current_tokens,
                "status": "fits" if current_tokens <= budget else "chunk_oversize",
                "extraction": {"paragraphs": range_spec(current_start, end)},
            }
        )
    return chunks


def plan_paragraph_file(path: Path, file_type: str, context_limit: int, reserve: int, item_budget: int) -> dict[str, Any]:
    paragraphs = docx_paragraphs(path) if file_type == "docx" else text_to_paragraphs(path)
    warnings: list[str] = []
    raw, structure_source, confidence = detect_paragraph_items(paragraphs, len(paragraphs), file_type)
    if len(raw) >= 2:
        items = make_paragraph_items(raw, paragraphs, item_budget)
        if any(item.get("status") == "chapter_oversize" for item in items):
            warnings.append("One or more heading sections exceeded the batch budget and were split.")
    else:
        items = make_auto_paragraph_chunks(paragraphs, item_budget)
        structure_source = "auto_chunks"
        confidence = "low"
        warnings.append("No reliable headings found; planned automatic paragraph/token chunks.")
    return make_plan(path, file_type, structure_source, confidence, None, items, context_limit, reserve, warnings)


def plan_pptx(path: Path, context_limit: int, reserve: int, item_budget: int) -> dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(str(path))
    slide_items = []
    for index, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                texts.append(shape.text.strip())
        title = compact_title(texts[0], 80) if texts else f"Slide {index}"
        slide_items.append({"slide": index, "title": title, "tokens": estimate_tokens("\n".join(texts))})

    chunks: list[dict[str, Any]] = []
    start = 1
    current = 0
    index = 1
    for slide in slide_items:
        if current and current + slide["tokens"] > item_budget:
            end = slide["slide"] - 1
            chunks.append(
                {
                    "id": f"chunk{index:03d}",
                    "title": f"幻灯片分块 {index}: slide {start}-{end}",
                    "kind": "auto_chunk",
                    "start_slide": start,
                    "end_slide": end,
                    "estimated_tokens": current,
                    "status": "fits",
                    "extraction": {"slides": range_spec(start, end)},
                }
            )
            index += 1
            start = slide["slide"]
            current = 0
        current += slide["tokens"]
    if slide_items:
        end = slide_items[-1]["slide"]
        chunks.append(
            {
                "id": f"chunk{index:03d}",
                "title": f"幻灯片分块 {index}: slide {start}-{end}",
                "kind": "auto_chunk",
                "start_slide": start,
                "end_slide": end,
                "estimated_tokens": current,
                "status": "fits" if current <= item_budget else "chunk_oversize",
                "extraction": {"slides": range_spec(start, end)},
            }
        )
    return make_plan(path, "pptx", "auto_chunks", "low", None, chunks, context_limit, reserve, ["PPTX is planned by slide title/continuous slide chunks."])


def make_batches(items: list[dict[str, Any]], batch_budget: int) -> list[dict[str, Any]]:
    batches: list[dict[str, Any]] = []
    current: list[dict[str, Any]] = []
    current_tokens = 0
    batch_index = 1
    for item in items:
        tokens = int(item.get("estimated_tokens") or 0)
        if current and current_tokens + tokens > batch_budget:
            batches.append(make_batch(batch_index, current, current_tokens))
            batch_index += 1
            current = []
            current_tokens = 0
        current.append(item)
        current_tokens += tokens
        if tokens > batch_budget:
            batches.append(make_batch(batch_index, current, current_tokens))
            batch_index += 1
            current = []
            current_tokens = 0
    if current:
        batches.append(make_batch(batch_index, current, current_tokens))
    return batches


def make_batch(index: int, items: list[dict[str, Any]], tokens: int) -> dict[str, Any]:
    return {
        "id": f"batch-{index}",
        "label": f"Batch {index}",
        "item_ids": [item["id"] for item in items],
        "items": items,
        "estimated_tokens": tokens,
        "status": "fits",
    }


def make_plan(
    path: Path,
    file_type: str,
    structure_source: str,
    confidence: str,
    page_offset: int | None,
    items: list[dict[str, Any]],
    context_limit: int,
    reserve: int,
    warnings: list[str],
) -> dict[str, Any]:
    available = max(1000, context_limit - reserve)
    batch_budget = max(1000, int(available * 0.75))
    batches = make_batches(items, batch_budget)
    return {
        "file": str(path),
        "file_name": path.name,
        "file_type": file_type,
        "context_limit": context_limit,
        "reserved_tokens": reserve,
        "available_tokens": available,
        "batch_budget": batch_budget,
        "structure_source": structure_source,
        "confidence": confidence,
        "page_offset": page_offset,
        "chapters": items,
        "all_items": [item.get("id") for item in items],
        "default_extraction": "all_items",
        "batches": batches,
        "first_batch": batches[0]["id"] if batches else None,
        "warnings": warnings,
    }


def render_markdown(plan: dict[str, Any]) -> str:
    lines = [
        "# 章节/分块处理计划",
        "",
        f"- 文件：`{plan.get('file_name')}`",
        f"- 类型：`{plan.get('file_type')}`",
        f"- 结构来源：`{plan.get('structure_source')}`",
        f"- 置信度：`{plan.get('confidence')}`",
        f"- 页码偏移：`{plan.get('page_offset')}`",
        f"- 可用 token：`{plan.get('available_tokens')}`",
        f"- 单元预算：`{plan.get('batch_budget')}`",
        f"- 默认抽取：`{plan.get('default_extraction', 'all_items')}`（全部内容单元）",
        f"- 兼容批次数：`{len(plan.get('batches', []))}`",
        "",
        "## 内容单元",
        "",
        "| ID | 标题 | 范围 | tokens | 状态 |",
        "| --- | --- | --- | ---: | --- |",
    ]
    for item in plan.get("chapters", []):
        extraction = item.get("extraction") or {}
        if "pages" in extraction:
            scope = f"p.{extraction['pages']}"
        elif "slides" in extraction:
            scope = f"slide {extraction['slides']}"
        else:
            scope = f"paragraph {extraction.get('paragraphs', '')}"
        lines.append(f"| {item.get('id')} | {item.get('title')} | {scope} | {item.get('estimated_tokens', 0)} | {item.get('status')} |")
    lines.extend(["", "## 兼容批次（非默认）", ""])
    for batch in plan.get("batches", []):
        titles = "；".join(item.get("title", "") for item in batch.get("items", []))
        lines.extend(
            [
                f"### {batch.get('id')}",
                "",
                f"- tokens：`{batch.get('estimated_tokens')}`",
                f"- 内容：{titles}",
                "",
            ]
        )
    warning_lines = "\n".join(f"- {warning}" for warning in plan.get("warnings", [])) or "- 无"
    lines.extend(["## 警告", "", warning_lines, ""])
    return "\n".join(lines)


def plan_document(path: Path, context_limit: int, reserve: int) -> dict[str, Any]:
    available = max(1000, context_limit - reserve)
    item_budget = max(1000, int(available * 0.75))
    ext = path.suffix.lower()
    if ext == ".pdf":
        return plan_pdf(path, context_limit, reserve, item_budget)
    if ext == ".docx":
        return plan_paragraph_file(path, "docx", context_limit, reserve, item_budget)
    if ext in {".md", ".markdown"}:
        return plan_paragraph_file(path, "markdown", context_limit, reserve, item_budget)
    if ext == ".txt":
        return plan_paragraph_file(path, "txt", context_limit, reserve, item_budget)
    if ext == ".pptx":
        return plan_pptx(path, context_limit, reserve, item_budget)
    raise SystemExit(f"Unsupported extension: {ext}")


def write_optional(path: str | None, content: str) -> None:
    if path:
        Path(path).write_text(content, encoding="utf-8")


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Plan chapters, sections, slide groups, or chunks for learning documents.")
    parser.add_argument("file")
    parser.add_argument("--context-limit", type=int, default=100000)
    parser.add_argument("--reserve", type=int, default=30000)
    parser.add_argument("--json-output")
    parser.add_argument("--markdown-output")
    parser.add_argument("--print-json", action="store_true")
    args = parser.parse_args(argv)

    plan = plan_document(Path(args.file), args.context_limit, args.reserve)
    json_text = json.dumps(plan, ensure_ascii=False, indent=2)
    markdown_text = render_markdown(plan)
    write_optional(args.json_output, json_text + "\n")
    write_optional(args.markdown_output, markdown_text)
    print(json_text if args.print_json else markdown_text)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
