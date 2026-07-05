#!/usr/bin/env python
"""Find visual candidates in PDF/PPTX learning materials without OCR."""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any


CAPTION_RE = re.compile(
    r"^\s*((?:fig(?:ure)?\.?|table|图|表)\s*[\dA-Za-z一二三四五六七八九十\-_.：:]*\s*.*)$",
    re.IGNORECASE,
)


def find_captions(text: str) -> list[str]:
    captions: list[str] = []
    for line in text.splitlines():
        match = CAPTION_RE.match(line.strip())
        if match:
            captions.append(match.group(1).strip())
    return captions[:5]


def nearby_text(text: str, limit: int = 500) -> str:
    compact = re.sub(r"\s+", " ", text).strip()
    return compact[:limit]


def classify_candidate(image_count: int, table_count: int, chart_count: int, vector_count: int, captions: list[str]) -> tuple[str, str]:
    if table_count:
        visual_type = "table"
    elif chart_count:
        visual_type = "chart"
    elif vector_count >= 20:
        visual_type = "diagram-or-vector-figure"
    elif image_count:
        visual_type = "image"
    else:
        visual_type = "text-referenced-visual"

    if captions or table_count or chart_count:
        importance = "high"
    elif image_count or vector_count >= 10:
        importance = "medium"
    else:
        importance = "low"
    return visual_type, importance


def audit_pdf_visuals(path: Path) -> list[dict[str, Any]]:
    import pdfplumber

    candidates: list[dict[str, Any]] = []
    with pdfplumber.open(str(path)) as pdf:
        for page_number, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            captions = find_captions(text)
            images = getattr(page, "images", None) or []
            vector_count = sum(
                len(getattr(page, attr, None) or [])
                for attr in ("lines", "rects", "curves")
            )
            try:
                tables = page.extract_tables() or []
            except Exception:
                tables = []
            image_count = len(images)
            table_count = len(tables)
            if not (image_count or table_count or vector_count >= 3 or captions):
                continue
            visual_type, importance = classify_candidate(image_count, table_count, 0, vector_count, captions)
            candidates.append(
                {
                    "source": f"{path.name} p.{page_number}",
                    "page": page_number,
                    "visual_type": visual_type,
                    "importance": importance,
                    "image_count": image_count,
                    "table_count": table_count,
                    "vector_shape_count": vector_count,
                    "caption_candidates": captions,
                    "nearby_text": nearby_text(text),
                    "uncertainty": "No OCR was used; text inside raster images may be unreadable.",
                    "suggested_note_action": "Create a simplified Mermaid diagram or ASCII line sketch only when nearby text/caption supports the interpretation.",
                }
            )
    return candidates


def shape_text(shape: Any) -> str:
    chunks: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip()
        if text:
            chunks.append(text)
    if getattr(shape, "has_table", False):
        try:
            for row in shape.table.rows:
                chunks.append(" | ".join(cell.text.strip() for cell in row.cells))
        except Exception:
            pass
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            child_text = shape_text(child)
            if child_text:
                chunks.append(child_text)
    return "\n".join(chunks)


def audit_pptx_visuals(path: Path) -> list[dict[str, Any]]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    candidates: list[dict[str, Any]] = []
    prs = Presentation(str(path))
    for slide_number, slide in enumerate(prs.slides, start=1):
        image_count = 0
        table_count = 0
        chart_count = 0
        vector_count = 0
        text_chunks: list[str] = []
        for shape in slide.shapes:
            text = shape_text(shape)
            if text:
                text_chunks.append(text)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
            if getattr(shape, "has_table", False):
                table_count += 1
            if getattr(shape, "has_chart", False):
                chart_count += 1
            shape_name = getattr(shape.shape_type, "name", str(shape.shape_type)).lower()
            if any(marker in shape_name for marker in ("auto_shape", "freeform", "line", "connector")):
                vector_count += 1
        all_text = "\n".join(text_chunks)
        captions = find_captions(all_text)
        if not (image_count or table_count or chart_count or vector_count >= 1 or captions):
            continue
        visual_type, importance = classify_candidate(image_count, table_count, chart_count, vector_count, captions)
        candidates.append(
            {
                "source": f"{path.name} slide {slide_number}",
                "slide": slide_number,
                "visual_type": visual_type,
                "importance": importance,
                "image_count": image_count,
                "table_count": table_count,
                "chart_count": chart_count,
                "vector_shape_count": vector_count,
                "caption_candidates": captions,
                "nearby_text": nearby_text(all_text),
                "uncertainty": "No OCR was used; text inside raster images may be unreadable.",
                "suggested_note_action": "Use slide text and captions to make a Mermaid or ASCII study sketch.",
            }
        )
    return candidates


def render_markdown(candidates: list[dict[str, Any]]) -> str:
    lines = ["# 重要图表候选", ""]
    if not candidates:
        lines.extend(["未发现明显图表候选。", ""])
        return "\n".join(lines)
    for index, item in enumerate(candidates, start=1):
        captions = item.get("caption_candidates") or []
        lines.extend(
            [
                f"## V{index}: {item.get('source')}",
                "",
                f"- 类型：`{item.get('visual_type')}`",
                f"- 重要性：`{item.get('importance')}`",
                f"- 图片数：`{item.get('image_count', 0)}`",
                f"- 表格数：`{item.get('table_count', 0)}`",
                f"- 图形/线条数：`{item.get('vector_shape_count', 0)}`",
                f"- 图注候选：{'; '.join(captions) if captions else '无'}",
                f"- 不确定性：{item.get('uncertainty')}",
                "",
                "附近文本：",
                "",
                f"> {item.get('nearby_text', '')}",
                "",
                "建议：",
                "",
                f"- {item.get('suggested_note_action')}",
                "",
            ]
        )
    return "\n".join(lines)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Detect important visual candidates in PDF/PPTX files without OCR.")
    parser.add_argument("file")
    parser.add_argument("--json-output", help="Optional JSON output path.")
    parser.add_argument("--markdown-output", help="Optional Markdown output path.")
    parser.add_argument("--print-json", action="store_true", help="Print JSON instead of Markdown.")
    args = parser.parse_args(argv)

    path = Path(args.file)
    if path.suffix.lower() == ".pdf":
        candidates = audit_pdf_visuals(path)
    elif path.suffix.lower() == ".pptx":
        candidates = audit_pptx_visuals(path)
    else:
        raise SystemExit("extract_visuals.py supports PDF and PPTX only.")

    payload = {"file": str(path), "candidates": candidates}
    json_text = json.dumps(payload, ensure_ascii=False, indent=2)
    markdown_text = render_markdown(candidates)
    if args.json_output:
        Path(args.json_output).write_text(json_text + "\n", encoding="utf-8")
    if args.markdown_output:
        Path(args.markdown_output).write_text(markdown_text, encoding="utf-8")
    print(json_text if args.print_json else markdown_text)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
